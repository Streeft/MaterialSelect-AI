"""PPTX renderer for Report (B2 — architecture, not yet a shipped export path).

Mirrors the CSV/XLSX renderers in spreadsheet.py: same Report input, same
"every export carries the limitation notice" guarantee (CLAUDE.md §1.8) —
here as its own slide rather than a header row, since a slide is the native
place for prose in this format. Not wired to a router endpoint; see
docs/TODO.md B2 and the module-level note in exporters/__init__.py.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt

from app.exporters.report import Report

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)


def to_pptx(report: Report) -> bytes:
    """Render ``report`` as a .pptx: title slide, one table slide per sheet,
    one closing slide with every notice."""
    presentation = Presentation()
    presentation.slide_width = _SLIDE_WIDTH
    presentation.slide_height = _SLIDE_HEIGHT

    _add_title_slide(presentation, report)
    for sheet in report.sheets:
        _add_table_slide(presentation, sheet.name, sheet.header, sheet.rows)
    if report.notices:
        _add_notices_slide(presentation, report.notices)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _add_title_slide(presentation: Presentation, report: Report) -> None:
    layout = presentation.slide_layouts[0]  # title layout
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = report.title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = report.subtitle


def _add_table_slide(
    presentation: Presentation, name: str, header: list[str], rows: list[list[object]]
) -> None:
    layout = presentation.slide_layouts[5]  # title-only layout
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = name

    n_rows = len(rows) + 1
    n_cols = max(len(header), 1)
    left, top = Inches(0.5), Inches(1.5)
    width, height = _SLIDE_WIDTH - Inches(1.0), _SLIDE_HEIGHT - Inches(2.0)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for col_index, label in enumerate(header):
        cell = table.cell(0, col_index)
        cell.text = str(label)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(14)

    for row_index, row in enumerate(rows, start=1):
        for col_index in range(n_cols):
            value = row[col_index] if col_index < len(row) else ""
            table.cell(row_index, col_index).text = "" if value is None else str(value)


def _add_notices_slide(presentation: Presentation, notices: list[str]) -> None:
    layout = presentation.slide_layouts[1]  # title + content
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = "Avisos"
    body = slide.placeholders[1].text_frame
    body.text = notices[0]
    for notice in notices[1:]:
        paragraph = body.add_paragraph()
        paragraph.text = notice
