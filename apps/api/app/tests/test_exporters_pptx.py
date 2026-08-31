from io import BytesIO

from pptx import Presentation

from app.exporters.pptx import to_pptx
from app.exporters.report import Report, Sheet, standard_notices


def test_to_pptx_produces_a_title_and_one_slide_per_sheet():
    report = Report(
        title="Seleção de materiais",
        subtitle="Estudo de exemplo",
        notices=standard_notices(includes_demo_data=True),
        sheets=[
            Sheet(
                name="Ranking",
                header=["Material", "Índice"],
                rows=[["Aço", 12.3], ["Alumínio", 9.8]],
            ),
        ],
    )
    data = to_pptx(report)
    assert isinstance(data, bytes)
    assert len(data) > 0

    presentation = Presentation(BytesIO(data))
    # Title slide + one slide per Sheet + one notices slide.
    assert len(list(presentation.slides)) == 3


def test_to_pptx_table_slide_has_header_and_data_rows():
    report = Report(
        title="T",
        subtitle="",
        notices=[],
        sheets=[Sheet(name="Dados", header=["A", "B"], rows=[[1, 2], [3, 4]])],
    )
    presentation = Presentation(BytesIO(to_pptx(report)))
    table_slide = list(presentation.slides)[1]
    tables = [shape for shape in table_slide.shapes if shape.has_table]
    assert len(tables) == 1
    table = tables[0].table
    assert table.cell(0, 0).text == "A"
    assert table.cell(0, 1).text == "B"
    assert table.cell(1, 0).text == "1"


def test_to_pptx_includes_the_limitation_notice():
    report = Report(
        title="T",
        subtitle="",
        notices=standard_notices(includes_demo_data=False),
        sheets=[Sheet(name="S", header=["x"], rows=[[1]])],
    )
    presentation = Presentation(BytesIO(to_pptx(report)))
    notices_slide_text = "\n".join(
        shape.text_frame.text
        for shape in list(presentation.slides)[-1].shapes
        if shape.has_text_frame
    )
    assert "didático" in notices_slide_text.lower()
